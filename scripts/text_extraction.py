import os
from pathlib import Path
import pymupdf        # pip install pymupdf
from pptx import Presentation
import nbformat

RAW_DIR = Path("corpus/raw")
OUT_DIR = Path("corpus/processed/texts")
OUT_DIR.mkdir(parents=True, exist_ok=True)

def extract_pdf(path: Path) -> str:
    """Extract text from a PDF using pymupdf."""
    text = []
    doc = pymupdf.open(str(path))
    for page in doc:
        text.append(page.get_text())
    doc.close()
    return "\n".join(text)

def extract_pptx(path: Path) -> str:
    """Extract all visible text from a PPTX file."""
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    return "\n".join(text)

def extract_notebook(path: Path) -> str:
    """Extract markdown and code cells from a .ipynb."""
    nb = nbformat.read(str(path), as_version=4)
    segments = []
    for cell in nb.cells:
        if cell.cell_type == "markdown":
            segments.append(cell.source)
        elif cell.cell_type == "code":
            segments.append("```python\n" + cell.source + "\n```")
    return "\n\n".join(segments)

def main():
    for root, _, files in os.walk(RAW_DIR):
        for fname in files:
            src = Path(root) / fname
            rel = src.relative_to(RAW_DIR)
            ext = src.suffix.lower()

            try:
                if ext == ".pdf":
                    content = extract_pdf(src)
                elif ext == ".pptx":
                    content = extract_pptx(src)
                elif ext == ".ipynb":
                    content = extract_notebook(src)
                else:
                    # skip other file types
                    continue

                dest = OUT_DIR / rel.with_suffix(".txt")
                dest.parent.mkdir(parents=True, exist_ok=True)
                dest.write_text(content, encoding="utf-8")
                print(f"Extracted {rel} → {dest}")

            except Exception as e:
                print(f"[Error] {rel}: {e}")

if __name__ == "__main__":
    main()
